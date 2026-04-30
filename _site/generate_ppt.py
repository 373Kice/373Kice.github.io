#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
A股大盘分析报告PPT生成器
基于2025年3月10日股市行情自动生成PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
import json

# 颜色定义
PRIMARY_COLOR = RGBColor(30, 58, 138)      # 深蓝色 #1E3A8A
SECONDARY_COLOR = RGBColor(59, 130, 246)   # 浅蓝色 #3B82F6
ACCENT_COLOR = RGBColor(16, 185, 129)      # 绿色 #10B981
DANGER_COLOR = RGBColor(239, 68, 68)       # 红色 #EF4444
TEXT_COLOR = RGBColor(31, 41, 55)           # 深灰色 #1F2937
LIGHT_TEXT_COLOR = RGBColor(107, 114, 128) # 浅灰色 #6B7280
BG_COLOR = RGBColor(255, 255, 255)         # 白色

# 字体大小
TITLE_FONT_SIZE = Pt(40)
SUBTITLE_FONT_SIZE = Pt(28)
HEADING_FONT_SIZE = Pt(32)
BODY_FONT_SIZE = Pt(22)
NOTE_FONT_SIZE = Pt(16)

class PPTGenerator:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)  # 1920px at 96 DPI
        self.prs.slide_height = Inches(7.5)    # 1080px at 96 DPI

    def add_title_slide(self):
        """添加封面页"""
        slide_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)

        # 标题
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "2025年3月10日A股大盘分析报告"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = TITLE_FONT_SIZE
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_COLOR
        title_para.alignment = PP_ALIGN.CENTER

        # 副标题
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11.333), Inches(0.6))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "缩量整理 结构分化 北向逆势流入"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = SUBTITLE_FONT_SIZE
        subtitle_para.font.color.rgb = SECONDARY_COLOR
        subtitle_para.alignment = PP_ALIGN.CENTER

        # 日期
        date_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.333), Inches(0.5))
        date_frame = date_box.text_frame
        date_frame.text = "2025年3月10日"
        date_para = date_frame.paragraphs[0]
        date_para.font.size = BODY_FONT_SIZE
        date_para.font.color.rgb = LIGHT_TEXT_COLOR
        date_para.alignment = PP_ALIGN.CENTER

        # 制作
        author_box = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(11.333), Inches(0.5))
        author_frame = author_box.text_frame
        author_frame.text = "制作: WorkBuddy AI"
        author_para = author_frame.paragraphs[0]
        author_para.font.size = NOTE_FONT_SIZE
        author_para.font.color.rgb = LIGHT_TEXT_COLOR
        author_para.alignment = PP_ALIGN.CENTER

    def add_content_slide(self, title, contents):
        """添加内容页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = HEADING_FONT_SIZE
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_COLOR

        # 内容
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.7))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        if isinstance(contents, list):
            for i, item in enumerate(contents):
                if i == 0:
                    p = content_frame.paragraphs[0]
                else:
                    p = content_frame.add_paragraph()
                p.text = item
                p.font.size = BODY_FONT_SIZE
                p.font.color.rgb = TEXT_COLOR
                p.space_after = Pt(10)
        elif isinstance(contents, str):
            content_frame.text = contents
            content_frame.paragraphs[0].font.size = BODY_FONT_SIZE
            content_frame.paragraphs[0].font.color.rgb = TEXT_COLOR

    def add_table_slide(self, title, headers, rows, highlight_positive=True):
        """添加表格页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = HEADING_FONT_SIZE
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_COLOR

        # 表格
        rows_count = len(rows) + 1  # 包括表头
        cols_count = len(headers)
        table = slide.shapes.add_table(rows_count, cols_count,
                                        Inches(0.5), Inches(1.3),
                                        Inches(12.333), Inches(5.7)).table

        # 表头
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.text_frame.paragraphs[0].font.size = Pt(20)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PRIMARY_COLOR

        # 数据行
        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_data in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(cell_data)
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(18)
                paragraph.font.color.rgb = TEXT_COLOR

                # 高亮正数
                if highlight_positive and isinstance(cell_data, str):
                    if "+" in cell_data or cell_data.endswith("%"):
                        # 检查是否是正数
                        try:
                            num = float(cell_data.replace("%", "").replace("+", "").replace(",", ""))
                            if num > 0:
                                paragraph.font.color.rgb = ACCENT_COLOR
                            elif num < 0:
                                paragraph.font.color.rgb = DANGER_COLOR
                        except:
                            pass

    def add_two_column_slide(self, title, left_content, right_content):
        """添加双栏内容页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = HEADING_FONT_SIZE
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_COLOR

        # 左栏
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(5.9), Inches(5.7))
        left_frame = left_box.text_frame
        left_frame.word_wrap = True

        for i, item in enumerate(left_content):
            if i == 0:
                p = left_frame.paragraphs[0]
            else:
                p = left_frame.add_paragraph()
            p.text = item
            p.font.size = BODY_FONT_SIZE
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(10)

        # 右栏
        right_box = slide.shapes.add_textbox(Inches(6.9), Inches(1.3), Inches(5.9), Inches(5.7))
        right_frame = right_box.text_frame
        right_frame.word_wrap = True

        for i, item in enumerate(right_content):
            if i == 0:
                p = right_frame.paragraphs[0]
            else:
                p = right_frame.add_paragraph()
            p.text = item
            p.font.size = BODY_FONT_SIZE
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(10)

    def add_list_slide(self, title, items, icon_prefix=""):
        """添加列表页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = HEADING_FONT_SIZE
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_COLOR

        # 列表
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.7))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        for i, item in enumerate(items):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()

            if icon_prefix:
                full_text = f"{icon_prefix} {item}"
            else:
                full_text = item

            p.text = full_text
            p.font.size = BODY_FONT_SIZE
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(15)
            p.line_spacing = 1.5

    def add_end_slide(self):
        """添加封底页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # 感谢文字
        thanks_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.333), Inches(1))
        thanks_frame = thanks_box.text_frame
        thanks_frame.text = "感谢观看"
        thanks_para = thanks_frame.paragraphs[0]
        thanks_para.font.size = TITLE_FONT_SIZE
        thanks_para.font.bold = True
        thanks_para.font.color.rgb = PRIMARY_COLOR
        thanks_para.alignment = PP_ALIGN.CENTER

        # 免责声明
        disclaimer_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(2))
        disclaimer_frame = disclaimer_box.text_frame
        disclaimer_frame.word_wrap = True

        disclaimer_text = """免责声明:
本报告基于公开信息整理,仅供参考,不构成任何投资建议。
股市有风险,投资需谨慎。
投资者据此操作,风险自担。

数据来源: 今日头条、新浪财经、雪球
更新时间: 2025年3月10日收盘"""

        disclaimer_frame.text = disclaimer_text
        for para in disclaimer_frame.paragraphs:
            para.font.size = NOTE_FONT_SIZE
            para.font.color.rgb = LIGHT_TEXT_COLOR
            para.alignment = PP_ALIGN.CENTER
            para.space_after = Pt(5)

    def generate(self):
        """生成完整的PPT"""

        # 1. 封面页
        self.add_title_slide()

        # 2. 目录页
        self.add_content_slide("目录", [
            "1. 大盘概况",
            "2. 资金流向",
            "3. 热点板块",
            "4. 个股资金流向TOP10",
            "5. 市场分析",
            "6. 后市展望",
            "7. 操作建议",
            "8. 风险提示",
            "9. 总结"
        ])

        # 3. 大盘概况 - 核心观点
        self.add_content_slide("大盘概况 - 核心观点", [
            "• 三大指数集体小幅收跌",
            "• 结束四连阳走势,进入震荡整理",
            "• 成交显著缩量,观望情绪升温",
            "• 个股涨多跌少,结构性机会突出",
            "",
            "关键数据:",
            "• 总成交额: 15057亿元 (↓3176亿元)",
            "• 涨跌分布: 3279只上涨 vs 1950只下跌",
            "• 上涨占比: 62.7%"
        ])

        # 4. 大盘概况 - 指数表现
        self.add_table_slide(
            "大盘概况 - 指数表现",
            ["指数", "收盘点位", "涨跌幅"],
            [
                ["上证指数", "3366.16点", "-0.19% ⬇️"],
                ["深证成指", "10825.70点", "-0.17% ⬇️"],
                ["创业板指", "2199.88点", "-0.25% ⬇️"],
                ["沪深300", "3928点", "-0.39% ⬇️"]
            ]
        )

        # 5. 资金流向 - 主力资金
        self.add_table_slide(
            "资金流向 - 主力资金",
            ["资金类型", "流向", "金额(亿元)"],
            [
                ["超大单", "净流出", "-167.27"],
                ["大单", "净流出", "-101.38"],
                ["中单", "净流入", "+58.26"],
                ["小单", "净流入", "+210.38"],
                ["合计", "净流出", "-248.13"]
            ]
        )

        # 6. 资金流向 - 北向资金
        self.add_content_slide("资金流向 - 北向资金", [
            "💰 净买入: +50亿元",
            "📊 逆势净流入",
            "",
            "外资态度分析:",
            "✅ 逆势净流入,显示外资看好",
            "✅ 与主力资金流出形成对比",
            "✅ 长期投资信心坚定"
        ])

        # 7. 热点板块 - 领涨板块
        self.add_table_slide(
            "热点板块 - 领涨板块TOP5",
            ["排名", "板块名称", "涨跌幅", "主力净流入(亿元)"],
            [
                ["🥇", "医药商业", "+4.08%", "+5.84"],
                ["🥈", "非金属材料", "+2.53%", "+6.79"],
                ["🥉", "电子元件", "+1.61%", "+6.68"],
                ["4", "煤炭", "+1.40%", "+4.64"],
                ["5", "有色金属", "+1.21%", "+7.85"]
            ]
        )

        # 8. 热点板块 - 领跌板块
        self.add_content_slide("热点板块 - 领跌板块", [
            "领跌板块TOP3:",
            "",
            "⚠️ 计算机: -1.63% (净流出-110.98亿)",
            "⚠️ 传媒: -0.96%",
            "⚠️ 非银金融: -0.84% (净流出-24.54亿)",
            "",
            "板块轮动逻辑:",
            "🔄 医药板块受政策利好推动",
            "🔄 资源板块表现强势",
            "🔄 科技板块调整明显",
            "🔄 符合市场轮动逻辑"
        ])

        # 9. 个股资金流向TOP10
        self.add_table_slide(
            "个股资金流向TOP10",
            ["排名", "股票名称", "主力净流入(亿元)", "涨跌幅"],
            [
                ["1", "洛阳钼业", "+9.55", "+8.62%"],
                ["2", "创业慧康", "+4.59", "+20.03%"],
                ["3", "胜宏科技", "+4.26", "+8.58%"],
                ["4", "寒武纪-U", "+4.25", "+3.87%"],
                ["5", "机器人", "+3.82", "+7.48%"],
                ["6", "海立股份", "+3.81", "+10.00%"],
                ["7", "大位科技", "+3.24", "+10.04%"],
                ["8", "海康威视", "+3.15", "+0.21%"],
                ["9", "三一重工", "+3.10", "+3.07%"],
                ["10", "东方锆业", "+3.02", "+10.00%"]
            ]
        )

        # 10. 市场分析 - 技术面
        self.add_two_column_slide(
            "市场分析 - 技术面",
            [
                "上证指数:",
                "• 当前点位: 3366.16点",
                "• 波动区间: 3350-3370点",
                "• 支撑位: 20日均线",
                "• 关键缺口: 3348点跳空缺口",
                "",
                "创业板指:",
                "• 当前点位: 2199.88点",
                "• 波动区间: 2180-2200点",
                "• 走势: 窄幅震荡",
                "• 状态: 缩量整理"
            ],
            [
                "走势判断:",
                "",
                "• 成长股经过前期上涨后",
                "  面临调整需求",
                "",
                "• 大盘蓝筹股在当前位置",
                "  承接力度较强",
                "",
                "• 整体维持震荡整理态势",
                "",
                "• 关注3350点支撑有效性"
            ]
        )

        # 11. 市场分析 - 市场情绪
        self.add_content_slide("市场分析 - 市场情绪", [
            "缩量调整特征:",
            "📉 成交额缩量至1.54万亿元",
            "📉 较上一日缩量3176亿元",
            "🎯 市场交投情绪趋于谨慎",
            "",
            "四连阳后的调整:",
            "🔍 技术性调整需求",
            "🔍 投资者观望情绪上升",
            "🔍 高位股分化明显",
            "",
            "市场情绪指标:",
            "• 恐慌指数: 适中",
            "• 贪婪指数: 中性",
            "• 观望指数: 偏高"
        ])

        # 12. 市场分析 - 市场特点
        self.add_list_slide(
            "市场分析 - 今日市场四大特点",
            [
                "1️⃣ 缩量整理 - 量能回落,观望升温",
                "2️⃣ 结构分化 - 板块轮动,机会不均",
                "3️⃣ 热点轮动 - 医药资源,科技调整",
                "4️⃣ 外资流入 - 北向逆势,长期看好"
            ],
            ""
        )

        # 13. 后市展望 - 短期趋势
        self.add_content_slide("后市展望 - 短期趋势", [
            "市场走势:",
            "📈 震荡整理为主",
            "📈 区间: 3350-3370点",
            "📈 方向: 横盘震荡",
            "",
            "关键因素:",
            "🔍 量能变化是关键",
            "🔍 关注支撑位有效性",
            "🔍 观望情绪能否改善",
            "",
            "操作建议:",
            "⏸️ 观望为主,谨慎操作",
            "⏸️ 控制仓位,不宜激进",
            "⏸️ 等待量能放大信号",
            "",
            "时间周期: 未来1-2周"
        ])

        # 14. 后市展望 - 中期趋势
        self.add_content_slide("后市展望 - 中期趋势", [
            "政策面:",
            "✅ 两会后财政发力预期",
            "✅ 适时降准降息政策",
            "✅ 产业政策支持",
            "",
            "基本面:",
            "✅ 经济复苏趋势确立",
            "✅ 企业盈利改善",
            "✅ 消费回暖",
            "",
            "流动性:",
            "✅ 流动性环境宽松",
            "✅ 资金面相对充裕",
            "✅ 融资环境改善",
            "",
            "市场走势: 稳步向上,结构性机会",
            "时间周期: 未来3-6个月"
        ])

        # 15. 后市展望 - 投资主线
        self.add_content_slide("后市展望 - 四大投资主线", [
            "🏥 主线一: 医药板块",
            "  政策受益明显,人口老龄化趋势",
            "",
            "🔋 主线二: 资源板块",
            "  周期性机会显现,大宗商品价格回升",
            "",
            "💻 主线三: 科技板块",
            "  调整后布局机会,AI技术持续突破",
            "",
            "🚀 主线四: 新质生产力",
            "  政策大力支持,新兴产业崛起"
        ])

        # 16. 操作建议 - 短期策略
        self.add_content_slide("操作建议 - 短期策略", [
            "仓位管理:",
            "• 总仓位: 控制在60%以内",
            "• 持仓周期: 1-2周",
            "• 调整频率: 不宜频繁",
            "",
            "具体建议:",
            "1. 控制仓位,不宜激进",
            "2. 关注3350-3370点支撑",
            "3. 把握结构性机会",
            "4. 等待量能放大信号",
            "",
            "操作要点:",
            "🎯 不追高  🎯 不杀跌",
            "🎯 逢高减仓  🎯 逢低吸纳"
        ])

        # 17. 操作建议 - 中期策略
        self.add_content_slide("操作建议 - 中期策略", [
            "核心原则:",
            "💎 精选优质标的",
            "💎 关注基本面",
            "💎 逢低布局",
            "💎 耐心持有",
            "",
            "选股标准:",
            "📈 业绩增长稳定",
            "📈 估值合理偏低",
            "📈 行业景气度高",
            "📈 竞争优势明显",
            "",
            "具体建议:",
            "1. 精选优质标的,关注基本面",
            "2. 关注政策受益板块",
            "3. 逢低布局核心资产",
            "4. 耐心持有,不频繁交易",
            "",
            "持仓周期: 3-6个月"
        ])

        # 18. 操作建议 - 重点关注
        self.add_content_slide("操作建议 - 四大重点关注", [
            "📊 量能变化",
            "   成交量是否放量,换手率是否活跃,资金流向是否改善",
            "",
            "📜 政策落地",
            "   财政政策是否发力,货币政策是否宽松,产业政策是否支持",
            "",
            "📈 行业景气",
            "   行业周期是否向上,供需关系是否改善,价格趋势是否向好",
            "",
            "💰 公司业绩",
            "   季报业绩是否超预期,订单是否增加,产能是否扩张",
            "",
            "跟踪频率: 每周复盘"
        ])

        # 19. 风险提示
        self.add_table_slide(
            "风险提示",
            ["风险类型", "风险描述", "应对措施"],
            [
                ["⚠️ 市场风险", "震荡调整可能持续", "控制仓位,分散投资"],
                ["⚠️ 政策风险", "政策落地不及预期", "关注政策动向"],
                ["⚠️ 流动性风险", "税期临近资金紧", "保持流动性"],
                ["⚠️ 外围风险", "海外市场波动", "关注外围市场"],
                ["⚠️ 个股风险", "高位股分化回调", "避免追高"]
            ]
        )

        # 20. 总结 - 核心观点
        self.add_content_slide("总结 - 核心观点", [
            "市场表现:",
            "✅ 三大指数小幅收跌",
            "✅ 震荡整理,观望升温",
            "✅ 成交显著缩量",
            "✅ 结构性分化明显",
            "",
            "资金特征:",
            "✅ 主力资金大幅流出",
            "✅ 北向资金逆势流入",
            "✅ 散户资金相对活跃",
            "✅ 外资长期看好",
            "",
            "市场机会:",
            "✅ 医药板块政策受益",
            "✅ 资源板块周期机会",
            "✅ 科技板块调整布局",
            "✅ 新质生产力长期方向"
        ])

        # 21. 总结 - 关键数据回顾
        self.add_content_slide("总结 - 关键数据回顾", [
            "指数表现:",
            "• 上证指数: 3366.16点 (↓0.19%)",
            "• 深证成指: 10825.70点 (↓0.17%)",
            "• 创业板指: 2199.88点 (↓0.25%)",
            "",
            "成交情况:",
            "• 成交额: 15057亿元 (↓3176亿元)",
            "• 涨跌分布: 3279涨 / 1950跌",
            "• 上涨占比: 62.7%",
            "",
            "资金流向:",
            "• 主力净流出: 248.13亿元",
            "• 北向净买入: 50亿元",
            "• 散户净流入: 210.38亿元",
            "",
            "领涨板块:",
            "• 医药商业: +4.08%",
            "• 非金属材料: +2.53%",
            "• 电子元件: +1.61%"
        ])

        # 22. 总结 - 投资建议
        self.add_content_slide("总结 - 投资建议", [
            "短期 (1-2周):",
            "⏸️ 控制仓位,谨慎观望",
            "⏸️ 不追高,不杀跌",
            "⏸️ 等待量能放大信号",
            "⏸️ 关注3350点支撑",
            "",
            "中期 (3-6个月):",
            "✅ 把握结构性机会",
            "✅ 逢低布局优质标的",
            "✅ 关注政策受益板块",
            "✅ 坚持价值投资",
            "",
            "长期 (1-3年):",
            "🎯 坚持价值投资理念",
            "🎯 关注核心资产",
            "🎯 布局新质生产力",
            "🎯 享受长期复利",
            "",
            "投资理念: 价值投资,长期持有,控制风险,稳健收益"
        ])

        # 23. 封底页
        self.add_end_slide()

        # 保存PPT
        output_path = "D:/github/373Kice.github.io/2025年3月10日A股大盘分析报告.pptx"
        self.prs.save(output_path)
        print(f"PPT has been generated: {output_path}")
        return output_path


if __name__ == "__main__":
    generator = PPTGenerator()
    output_path = generator.generate()
    print(f"\n[Success] PPT generated successfully!")
    print(f"[File Path] {output_path}")
    print(f"[Total Pages] {generator.prs.slides.__len__()} pages")
