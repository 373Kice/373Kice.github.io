#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
A股大盘汇报PPT生成器（专业美观版）
深色商务风格，含渐变背景、色块装饰、数据可视化
"""

import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from lxml import etree
import copy

# ───────────────────────── 颜色体系 ─────────────────────────
C_BG        = RGBColor(0x0D, 0x1B, 0x3E)   # 深蓝背景
C_BG2       = RGBColor(0x13, 0x26, 0x52)   # 次深蓝
C_ACCENT    = RGBColor(0x00, 0xB0, 0xFF)   # 亮蓝强调
C_GREEN     = RGBColor(0x00, 0xE6, 0x76)   # 上涨绿
C_RED       = RGBColor(0xF4, 0x43, 0x36)   # 下跌红
C_ORANGE    = RGBColor(0xFF, 0xA7, 0x26)   # 橙色警示
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY      = RGBColor(0xB0, 0xBE, 0xC5)
C_CARD      = RGBColor(0x1A, 0x30, 0x60)   # 卡片背景
C_HEADER    = RGBColor(0x00, 0x5C, 0x9E)   # 表头蓝

# ───────────────────────── 幻灯片尺寸 ─────────────────────────
W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

def blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

# ───────── 通用辅助函数 ─────────

def fill_bg(slide, color=C_BG):
    """填充幻灯片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, color, alpha=None):
    """添加纯色矩形色块"""
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE.RECTANGLE=1
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, text, x, y, w, h, font_size=Pt(18),
                bold=False, color=C_WHITE, align=PP_ALIGN.LEFT,
                wrap=True, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb

def add_tag_rect(slide, label, x, y, w=Inches(2.2), h=Inches(0.4), bg=C_ACCENT):
    """小标签色块"""
    add_rect(slide, x, y, w, h, bg)
    add_textbox(slide, label, x + Inches(0.1), y + Inches(0.02),
                w - Inches(0.2), h, font_size=Pt(14), bold=True,
                color=C_WHITE, align=PP_ALIGN.CENTER)

def add_card(slide, title, lines, x, y, w, h,
             title_bg=C_HEADER, card_bg=C_CARD,
             title_color=C_WHITE, body_color=C_WHITE,
             title_fsize=Pt(18), body_fsize=Pt(15)):
    """带标题栏的卡片"""
    # 卡片底色
    add_rect(slide, x, y, w, h, card_bg)
    # 标题条
    th = Inches(0.45)
    add_rect(slide, x, y, w, th, title_bg)
    add_textbox(slide, title, x + Inches(0.15), y + Inches(0.05),
                w - Inches(0.3), th, font_size=title_fsize,
                bold=True, color=title_color, align=PP_ALIGN.LEFT)
    # 内容
    tb = slide.shapes.add_textbox(x + Inches(0.18), y + th + Inches(0.1),
                                   w - Inches(0.36), h - th - Inches(0.1))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.space_after = Pt(4)
        if p.runs:
            p.runs[0].font.size = body_fsize
            p.runs[0].font.color.rgb = body_color
        else:
            run = p.add_run()
            run.text = line
            run.font.size = body_fsize
            run.font.color.rgb = body_color

def add_table(slide, headers, rows, x, y, w, h,
              header_bg=C_HEADER, row_alt=C_CARD, row_bg=C_BG2,
              header_fsize=Pt(16), body_fsize=Pt(14),
              col_colors=None):
    """添加带样式的表格"""
    nrows = len(rows) + 1
    ncols = len(headers)
    tbl = slide.shapes.add_table(nrows, ncols, x, y, w, h).table

    # 列宽平分
    col_w = w // ncols
    for i in range(ncols):
        tbl.columns[i].width = col_w

    def set_cell(cell, text, bg, fg, fsize, bold=False, align=PP_ALIGN.CENTER):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size = fsize
        run.font.bold = bold
        run.font.color.rgb = fg

    # 表头
    for ci, h_text in enumerate(headers):
        set_cell(tbl.cell(0, ci), h_text, header_bg, C_WHITE, header_fsize, bold=True)

    # 数据行
    for ri, row in enumerate(rows):
        bg = row_alt if ri % 2 == 0 else row_bg
        for ci, val in enumerate(row):
            fg = C_WHITE
            # 自动上涨绿/下跌红
            s = str(val)
            if col_colors and ci < len(col_colors) and col_colors[ci]:
                fg = col_colors[ci]
            elif ('+' in s and '%' in s) or ('亿' in s and '+' in s):
                fg = C_GREEN
            elif ('-' in s and '%' in s and '↓' not in s) or ('亿' in s and '-' in s):
                fg = C_RED
            elif '↓' in s:
                fg = C_RED
            elif '↑' in s:
                fg = C_GREEN
            set_cell(tbl.cell(ri + 1, ci), val, bg, fg, body_fsize)
    return tbl

def slide_header(slide, title, subtitle=None):
    """添加统一顶部标题栏"""
    add_rect(slide, 0, 0, W, Inches(0.85), C_HEADER)
    add_textbox(slide, title,
                Inches(0.4), Inches(0.08), Inches(10), Inches(0.65),
                font_size=Pt(26), bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(10.5), Inches(0.2), Inches(2.6), Inches(0.5),
                    font_size=Pt(14), color=C_ACCENT, align=PP_ALIGN.RIGHT)
    # 底部装饰线
    add_rect(slide, 0, Inches(7.35), W, Inches(0.15), C_ACCENT)

# ═══════════════════════════════════════════════
# Slide 1 — 封面
# ═══════════════════════════════════════════════
def slide_cover():
    slide = blank_slide()
    fill_bg(slide, C_BG)

    # 左侧深色条
    add_rect(slide, 0, 0, Inches(0.2), H, C_ACCENT)

    # 右上角装饰三角（用细长矩形模拟）
    add_rect(slide, Inches(10.5), 0, Inches(2.833), Inches(2.0), C_BG2)
    add_rect(slide, Inches(12), 0, Inches(1.333), Inches(3.5), C_CARD)

    # 顶部标签
    add_tag_rect(slide, "2025年A股市场分析",
                 Inches(0.8), Inches(1.2), Inches(3.2), Inches(0.45))

    # 主标题
    add_textbox(slide, "A股大盘",
                Inches(0.8), Inches(2.0), Inches(9), Inches(1.2),
                font_size=Pt(56), bold=True, color=C_WHITE)
    add_textbox(slide, "分析报告",
                Inches(0.8), Inches(3.0), Inches(9), Inches(1.2),
                font_size=Pt(56), bold=True, color=C_ACCENT)

    # 副标题
    add_textbox(slide, "缩量整理  |  结构分化  |  北向逆势流入",
                Inches(0.8), Inches(4.3), Inches(9), Inches(0.6),
                font_size=Pt(22), color=C_GRAY)

    # 底部信息栏
    add_rect(slide, 0, Inches(6.5), W, Inches(1.0), C_CARD)
    add_textbox(slide, "报告日期：2025年3月10日",
                Inches(0.8), Inches(6.6), Inches(4), Inches(0.5),
                font_size=Pt(16), color=C_GRAY)
    add_textbox(slide, "数据来源：新浪财经 / 东方财富 / 雪球",
                Inches(5), Inches(6.6), Inches(5), Inches(0.5),
                font_size=Pt(16), color=C_GRAY, align=PP_ALIGN.CENTER)
    add_textbox(slide, "WorkBuddy AI",
                Inches(10.5), Inches(6.6), Inches(2.5), Inches(0.5),
                font_size=Pt(16), color=C_ACCENT, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════
# Slide 2 — 目录
# ═══════════════════════════════════════════════
def slide_toc():
    slide = blank_slide()
    fill_bg(slide, C_BG)
    slide_header(slide, "目录  CONTENTS", "2025.03.10")

    chapters = [
        ("01", "大盘概况", "三大指数 · 成交量 · 涨跌分布"),
        ("02", "资金流向", "主力资金 · 北向资金"),
        ("03", "热点板块", "领涨板块 · 领跌板块 · 轮动逻辑"),
        ("04", "个股TOP10", "主力净流入前十个股"),
        ("05", "市场分析", "技术面 · 情绪面 · 四大特点"),
        ("06", "后市展望", "短期 · 中期 · 四大投资主线"),
        ("07", "操作建议", "短期策略 · 中期策略 · 重点关注"),
        ("08", "风险提示", "五大风险 · 应对措施"),
        ("09", "总结", "核心观点 · 数据回顾 · 投资建议"),
    ]

    col_w  = Inches(4.1)
    row_h  = Inches(0.65)
    start_x = [Inches(0.4), Inches(4.7), Inches(9.0)]
    start_y  = Inches(1.1)

    for i, (num, title, sub) in enumerate(chapters):
        col = i % 3
        row = i // 3
        x = start_x[col]
        y = start_y + row * (row_h + Inches(0.18))

        # 序号方块
        add_rect(slide, x, y, Inches(0.55), row_h, C_ACCENT)
        add_textbox(slide, num, x, y, Inches(0.55), row_h,
                    font_size=Pt(16), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        # 内容块
        add_rect(slide, x + Inches(0.55), y, col_w - Inches(0.55), row_h, C_CARD)
        add_textbox(slide, title,
                    x + Inches(0.7), y + Inches(0.04), col_w - Inches(0.8), Inches(0.32),
                    font_size=Pt(16), bold=True, color=C_WHITE)
        add_textbox(slide, sub,
                    x + Inches(0.7), y + Inches(0.33), col_w - Inches(0.8), Inches(0.28),
                    font_size=Pt(11), color=C_GRAY)

# ═══════════════════════════════════════════════
# Slide 3 — 大盘概况（数字看板）
# ═══════════════════════════════════════════════
def slide_overview():
    slide = blank_slide()
    fill_bg(slide, C_BG)
    slide_header(slide, "01  大盘概况", "指数表现")

    # 核心观点标签
    add_tag_rect(slide, "今日市场概况",
                 Inches(0.4), Inches(1.0), Inches(2.4), Inches(0.38), C_ACCENT)

    add_textbox(slide,
        "三大指数集体小幅收跌，结束四连阳走势进入震荡整理；成交显著缩量，观望情绪升温；个股涨多跌少，结构性机会突出。",
        Inches(0.4), Inches(1.5), Inches(7.5), Inches(0.8),
        font_size=Pt(15), color=C_GRAY, wrap=True)

    # 指数数字卡片
    index_data = [
        ("上证指数", "3366.16", "-0.19%", "↓"),
        ("深证成指", "10825.70", "-0.17%", "↓"),
        ("创业板指", "2199.88",  "-0.25%", "↓"),
        ("沪深300",  "3928.00",  "-0.39%", "↓"),
    ]
    card_w = Inches(2.9)
    card_h = Inches(1.8)
    card_y = Inches(2.4)
    for i, (name, point, chg, arrow) in enumerate(index_data):
        cx = Inches(0.4) + i * (card_w + Inches(0.3))
        add_rect(slide, cx, card_y, card_w, card_h, C_CARD)
        add_rect(slide, cx, card_y, card_w, Inches(0.06), C_RED)  # 顶部红线
        add_textbox(slide, name, cx, card_y + Inches(0.15), card_w, Inches(0.4),
                    font_size=Pt(15), color=C_GRAY, align=PP_ALIGN.CENTER)
        add_textbox(slide, point, cx, card_y + Inches(0.55), card_w, Inches(0.7),
                    font_size=Pt(28), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, f"{arrow}  {chg}",
                    cx, card_y + Inches(1.25), card_w, Inches(0.45),
                    font_size=Pt(20), bold=True, color=C_RED, align=PP_ALIGN.CENTER)

    # 成交量 / 涨跌分布 统计卡
    stat_y = Inches(4.4)
    stats = [
        ("总成交额", "15057亿元", C_ACCENT),
        ("较前日缩量", "▼ 3176亿元", C_ORANGE),
        ("上涨家数", "3279只", C_GREEN),
        ("下跌家数", "1950只", C_RED),
        ("上涨占比", "62.7%", C_GREEN),
    ]
    sw = Inches(2.4)
    sh = Inches(1.3)
    for i, (label, val, color) in enumerate(stats):
        sx = Inches(0.4) + i * (sw + Inches(0.25))
        add_rect(slide, sx, stat_y, sw, sh, C_CARD)
        add_rect(slide, sx, stat_y + sh - Inches(0.06), sw, Inches(0.06), color)
        add_textbox(slide, label, sx, stat_y + Inches(0.08), sw, Inches(0.38),
                    font_size=Pt(13), color=C_GRAY, align=PP_ALIGN.CENTER)
        add_textbox(slide, val,   sx, stat_y + Inches(0.46), sw, Inches(0.6),
                    font_size=Pt(20), bold=True, color=color, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# Slide 4 — 资金流向
# ═══════════════════════════════════════════════
def slide_capital():
    slide = blank_slide()
    fill_bg(slide, C_BG)
    slide_header(slide, "02  资金流向", "主力资金 · 北向资金")

    # 左：主力资金表格
    add_tag_rect(slide, "主力资金", Inches(0.4), Inches(1.0),
                 Inches(2), Inches(0.38), C_ACCENT)
    add_table(slide,
              ["资金类型", "流向", "金额(亿元)"],
              [
                  ["超大单", "净流出", "-167.27"],
                  ["大单",   "净流出", "-101.38"],
                  ["中单",   "净流入",  "+58.26"],
                  ["小单",   "净流入", "+210.38"],
                  ["合计",   "净流出", "-248.13"],
              ],
              Inches(0.4), Inches(1.5), Inches(5.8), Inches(3.6))

    # 右上：北向资金突出数字
    add_tag_rect(slide, "北向资金（逆势净买入）",
                 Inches(6.8), Inches(1.0), Inches(3.5), Inches(0.38), C_GREEN)
    add_rect(slide, Inches(6.8), Inches(1.5), Inches(6.2), Inches(1.8), C_CARD)
    add_rect(slide, Inches(6.8), Inches(1.5), Inches(6.2), Inches(0.06), C_GREEN)
    add_textbox(slide, "+50亿元",
                Inches(6.8), Inches(1.7), Inches(6.2), Inches(1.0),
                font_size=Pt(52), bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
    add_textbox(slide, "外资逆势净流入，长期信心坚定",
                Inches(6.8), Inches(2.7), Inches(6.2), Inches(0.45),
                font_size=Pt(16), color=C_GRAY, align=PP_ALIGN.CENTER)

    # 右下：对比说明
    notes = [
        "主力资金大幅净流出 248亿，超大单主导流出",
        "北向资金逆势净买入 50亿，与主力形成反差",
        "散户（小单）净流入 210亿，参与意愿较强",
        "资金结构分化，外资长期配置逻辑未变",
    ]
    add_tag_rect(slide, "资金解读", Inches(6.8), Inches(3.5),
                 Inches(2), Inches(0.38), C_ORANGE)
    add_card(slide, "核心要点", notes,
             Inches(6.8), Inches(3.95), Inches(6.2), Inches(2.8),
             title_bg=C_CARD, card_bg=C_BG2,
             title_color=C_ACCENT, body_color=C_GRAY, body_fsize=Pt(14))

    # 底部进度条示意（主力 vs 北向）
    add_rect(slide, Inches(0.4), Inches(5.3), Inches(5.8), Inches(0.08), C_CARD)
    add_rect(slide, Inches(0.4), Inches(5.3), Inches(5.8 * 248/300), Inches(0.08), C_RED)
    add_textbox(slide, "主力净流出  248.13亿",
                Inches(0.4), Inches(5.4), Inches(5.8), Inches(0.4),
                font_size=Pt(13), color=C_GRAY)

    add_rect(slide, Inches(0.4), Inches(5.9), Inches(5.8), Inches(0.08), C_CARD)
    add_rect(slide, Inches(0.4), Inches(5.9), Inches(5.8 * 50/300), Inches(0.08), C_GREEN)
    add_textbox(slide, "北向净买入  50亿",
                Inches(0.4), Inches(6.0), Inches(5.8), Inches(0.4),
                font_size=Pt(13), color=C_GRAY)

# ═══════════════════════════════════════════════
# Slide 5 — 热点板块
# ═══════════════════════════════════════════════
def slide_sectors():
    slide = blank_slide()
    fill_bg(slide, C_BG)
    slide_header(slide, "03  热点板块", "领涨 · 领跌 · 轮动逻辑")

    # 领涨板块表格（左）
    add_tag_rect(slide, "领涨板块 TOP5",
                 Inches(0.4), Inches(1.0), Inches(2.4), Inches(0.38), C_GREEN)
    add_table(slide,
              ["板块", "涨跌幅", "主力净流入"],
              [
                  ["医药商业",   "+4.08%", "+5.84亿"],
                  ["非金属材料", "+2.53%", "+6.79亿"],
                  ["电子元件",   "+1.61%", "+6.68亿"],
                  ["煤炭",       "+1.40%", "+4.64亿"],
                  ["有色金属",   "+1.21%", "+7.85亿"],
              ],
              Inches(0.4), Inches(1.5), Inches(6.0), Inches(3.3))

    # 领跌板块表格（右）
    add_tag_rect(slide, "领跌板块 TOP3",
                 Inches(7.0), Inches(1.0), Inches(2.4), Inches(0.38), C_RED)
    add_table(slide,
              ["板块", "涨跌幅", "主力净流出"],
              [
                  ["计算机",  "-1.63%", "-110.98亿"],
                  ["传媒",    "-0.96%", "-"],
                  ["非银金融","-0.84%", "-24.54亿"],
              ],
              Inches(7.0), Inches(1.5), Inches(6.0), Inches(2.1))

    # 板块轮动逻辑卡片
    add_tag_rect(slide, "板块轮动逻辑",
                 Inches(7.0), Inches(3.8), Inches(2.6), Inches(0.38), C_ORANGE)
    add_card(slide, "今日轮动特征", [
        "医药板块受政策利好推动，强势领涨",
        "资源板块（有色、煤炭）表现强势",
        "计算机科技板块调整，净流出巨大",
        "板块轮动符合市场规律，结构性机会突出",
    ], Inches(7.0), Inches(4.25), Inches(6.0), Inches(2.3),
       title_bg=C_CARD, card_bg=C_BG2,
       title_color=C_ORANGE, body_color=C_GRAY, body_fsize=Pt(14))

    # 进度条可视化板块涨幅
    bars = [
        ("医药商业",   4.08, C_GREEN),
        ("非金属材料", 2.53, C_GREEN),
        ("电子元件",   1.61, C_ACCENT),
        ("煤炭",       1.40, C_ACCENT),
        ("有色金属",   1.21, C_ACCENT),
    ]
    bar_y = Inches(5.05)
    for i, (name, pct, color) in enumerate(bars):
        bx = Inches(0.4)
        by = bar_y + i * Inches(0.34)
        add_textbox(slide, name, bx, by, Inches(1.8), Inches(0.3),
                    font_size=Pt(12), color=C_GRAY, align=PP_ALIGN.RIGHT)
        add_rect(slide, bx + Inches(1.85), by + Inches(0.06), Inches(4.0), Inches(0.2), C_CARD)
        bar_len = Inches(4.0 * pct / 5.0)
        add_rect(slide, bx + Inches(1.85), by + Inches(0.06), bar_len, Inches(0.2), color)
        add_textbox(slide, f"+{pct}%",
                    bx + Inches(1.85) + bar_len + Inches(0.05), by, Inches(0.6), Inches(0.3),
                    font_size=Pt(12), color=color)

# ═══════════════════════════════════════════════
# Slide 6 — 个股资金流向TOP10
# ═══════════════════════════════════════════════
def slide_top10():
    slide = blank_slide()
    fill_bg(slide, C_BG)
    slide_header(slide, "04  个股资金流向 TOP10", "主力净流入最多个股")

    add_table(slide,
              ["排名", "股票名称", "主力净流入(亿元)", "今日涨跌幅"],
              [
                  ["1",  "洛阳钼业",  "+9.55", "+8.62%"],
                  ["2",  "创业慧康",  "+4.59", "+20.03%"],
                  ["3",  "胜宏科技",  "+4.26", "+8.58%"],
                  ["4",  "寒武纪-U",  "+4.25", "+3.87%"],
                  ["5",  "机器人",    "+3.82", "+7.48%"],
                  ["6",  "海立股份",  "+3.81", "+10.00%"],
                  ["7",  "大位科技",  "+3.24", "+10.04%"],
                  ["8",  "海康威视",  "+3.15", "+0.21%"],
                  ["9",  "三一重工",  "+3.10", "+3.07%"],
                  ["10", "东方锆业",  "+3.02", "+10.00%"],
              ],
              Inches(0.4), Inches(1.0), Inches(8.0), Inches(6.0))

    # 右侧亮点
    highlights = [
        ("洛阳钼业", "9.55亿", "资源板块龙头", C_GREEN),
        ("创业慧康", "4.59亿", "医药+20%强势", C_ACCENT),
        ("寒武纪-U", "4.25亿", "AI芯片长期主线", C_ORANGE),
        ("机器人",   "3.82亿", "机器人概念活跃", C_ACCENT),
    ]
    hx = Inches(8.7)
    hy = Inches(1.0)
    for name, val, desc, color in highlights:
        add_rect(slide, hx, hy, Inches(4.3), Inches(1.28), C_CARD)
        add_rect(slide, hx, hy, Inches(0.06), Inches(1.28), color)
        add_textbox(slide, name, hx + Inches(0.2), hy + Inches(0.12),
                    Inches(2.2), Inches(0.42), font_size=Pt(17), bold=True, color=C_WHITE)
        add_textbox(slide, val,  hx + Inches(2.5), hy + Inches(0.12),
                    Inches(1.6), Inches(0.42), font_size=Pt(17), bold=True, color=color, align=PP_ALIGN.RIGHT)
        add_textbox(slide, desc, hx + Inches(0.2), hy + Inches(0.6),
                    Inches(3.9), Inches(0.42), font_size=Pt(13), color=C_GRAY)
        hy += Inches(1.38)

# ═══════════════════════════════════════════════
# Slide 7 — 市场分析
# ═══════════════════════════════════════════════
def slide_analysis():
    slide = blank_slide()
    fill_bg(slide, C_BG)
    slide_header(slide, "05  市场分析", "技术面 · 情绪面 · 四大特点")

    # 技术面（左卡片）
    add_card(slide, "技术面分析", [
        "上证指数 3366 点，区间 3350-3370",
        "支撑位：20日均线 & 3348缺口",
        "创业板 2200 点附近窄幅震荡",
        "成长股调整需求，蓝筹承接较强",
        "整体维持震荡整理，关注3350支撑",
    ], Inches(0.4), Inches(1.0), Inches(6.0), Inches(2.8),
       title_bg=C_HEADER, card_bg=C_CARD,
       title_color=C_ACCENT, body_color=C_GRAY, body_fsize=Pt(15))

    # 市场情绪（右卡片）
    add_card(slide, "市场情绪分析", [
        "成交额缩量至 1.54万亿 (-3176亿)",
        "四连阳后进入技术性调整",
        "投资者观望情绪上升",
        "恐慌指数：适中 | 观望指数：偏高",
        "高位股分化，云鼎科技等跌停",
    ], Inches(6.8), Inches(1.0), Inches(6.2), Inches(2.8),
       title_bg=C_HEADER, card_bg=C_CARD,
       title_color=C_ORANGE, body_color=C_GRAY, body_fsize=Pt(15))

    # 四大特点大卡片
    add_tag_rect(slide, "今日市场四大特点",
                 Inches(0.4), Inches(4.0), Inches(2.8), Inches(0.4), C_ACCENT)

    features = [
        ("01", "缩量整理", "量能回落，观望升温",   C_RED),
        ("02", "结构分化", "板块轮动，机会不均",   C_ORANGE),
        ("03", "热点轮动", "医药资源强，科技调整", C_GREEN),
        ("04", "外资流入", "北向逆势，长期看好",   C_ACCENT),
    ]
    fw = Inches(2.9)
    fh = Inches(2.3)
    for i, (num, title, sub, color) in enumerate(features):
        fx = Inches(0.4) + i * (fw + Inches(0.4))
        fy = Inches(4.5)
        add_rect(slide, fx, fy, fw, fh, C_CARD)
        add_rect(slide, fx, fy, fw, Inches(0.08), color)
        add_textbox(slide, num, fx, fy + Inches(0.12), Inches(0.7), Inches(0.5),
                    font_size=Pt(22), bold=True, color=color, align=PP_ALIGN.CENTER)
        add_textbox(slide, title, fx + Inches(0.65), fy + Inches(0.12),
                    fw - Inches(0.7), Inches(0.5),
                    font_size=Pt(20), bold=True, color=C_WHITE)
        add_textbox(slide, sub, fx + Inches(0.15), fy + Inches(0.7),
                    fw - Inches(0.3), Inches(0.5),
                    font_size=Pt(15), color=C_GRAY)

# ═══════════════════════════════════════════════
# Slide 8 — 后市展望
# ═══════════════════════════════════════════════
def slide_outlook():
    slide = blank_slide()
    fill_bg(slide, C_BG)
    slide_header(slide, "06  后市展望", "短期 · 中期 · 投资主线")

    # 短期（左上）
    add_card(slide, "短期趋势（1-2周）", [
        "震荡整理为主，3350-3370区间",
        "量能变化是关键信号",
        "关注3348跳空缺口支撑",
        "策略：观望为主，控制仓位",
    ], Inches(0.4), Inches(1.0), Inches(6.0), Inches(2.5),
       title_bg=C_ORANGE, card_bg=C_CARD,
       title_color=C_WHITE, body_color=C_GRAY, body_fsize=Pt(15))

    # 中期（右上）
    add_card(slide, "中期趋势（3-6个月）", [
        "政策面：两会后财政发力，降准降息预期",
        "基本面：经济复苏，企业盈利改善",
        "流动性：环境宽松，资金面充裕",
        "市场走势：稳步向上，结构性机会",
    ], Inches(6.8), Inches(1.0), Inches(6.2), Inches(2.5),
       title_bg=C_GREEN, card_bg=C_CARD,
       title_color=C_WHITE, body_color=C_GRAY, body_fsize=Pt(15))

    # 四大投资主线
    add_tag_rect(slide, "四大投资主线",
                 Inches(0.4), Inches(3.7), Inches(2.4), Inches(0.4), C_ACCENT)

    mainlines = [
        ("医药板块", "政策受益+人口老龄化", "关注CXO、创新药、医疗器械",   C_ACCENT, "01"),
        ("资源板块", "周期机会+供给侧改革", "关注有色金属、煤炭、石油",     C_GREEN,  "02"),
        ("科技板块", "调整后布局+AI突破",   "关注AI芯片、算力、软件",       C_ORANGE, "03"),
        ("新质生产力","政策重点支持方向",    "关注机器人、低空经济、新材料", C_RED,    "04"),
    ]
    mw = Inches(2.9)
    mh = Inches(2.6)
    for i, (title, sub1, sub2, color, num) in enumerate(mainlines):
        mx = Inches(0.4) + i * (mw + Inches(0.4))
        my = Inches(4.2)
        add_rect(slide, mx, my, mw, mh, C_CARD)
        # 顶部彩色条
        add_rect(slide, mx, my, mw, Inches(0.07), color)
        # 序号圆角色块模拟
        add_rect(slide, mx + Inches(0.15), my + Inches(0.15),
                 Inches(0.45), Inches(0.45), color)
        add_textbox(slide, num, mx + Inches(0.15), my + Inches(0.12),
                    Inches(0.45), Inches(0.45),
                    font_size=Pt(13), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, title,
                    mx + Inches(0.7), my + Inches(0.18), mw - Inches(0.8), Inches(0.42),
                    font_size=Pt(17), bold=True, color=C_WHITE)
        add_textbox(slide, sub1,
                    mx + Inches(0.15), my + Inches(0.72), mw - Inches(0.3), Inches(0.4),
                    font_size=Pt(13), color=color)
        add_textbox(slide, sub2,
                    mx + Inches(0.15), my + Inches(1.15), mw - Inches(0.3), Inches(0.7),
                    font_size=Pt(12), color=C_GRAY, wrap=True)

# ═══════════════════════════════════════════════
# Slide 9 — 操作建议
# ═══════════════════════════════════════════════
def slide_strategy():
    slide = blank_slide()
    fill_bg(slide, C_BG)
    slide_header(slide, "07  操作建议", "短期 · 中期 · 长期策略")

    # 三列策略卡片
    strategies = [
        ("短期策略\n（1-2周）", [
            "仓位控制在 60% 以内",
            "不追高，不杀跌",
            "关注 3350 点支撑",
            "把握结构性机会",
            "等待量能放大信号",
        ], C_ORANGE),
        ("中期策略\n（3-6个月）", [
            "精选优质标的",
            "关注基本面改善",
            "逢低布局核心资产",
            "耐心持有不频繁操作",
            "持仓周期：3-6个月",
        ], C_GREEN),
        ("长期策略\n（1-3年）", [
            "坚持价值投资理念",
            "关注核心资产",
            "布局新质生产力",
            "分散配置降风险",
            "享受长期复利增长",
        ], C_ACCENT),
    ]
    sw = Inches(3.9)
    for i, (title, points, color) in enumerate(strategies):
        sx = Inches(0.4) + i * (sw + Inches(0.3))
        add_rect(slide, sx, Inches(1.0), sw, Inches(4.8), C_CARD)
        add_rect(slide, sx, Inches(1.0), sw, Inches(0.07), color)
        add_textbox(slide, title, sx, Inches(1.08), sw, Inches(0.9),
                    font_size=Pt(18), bold=True, color=color, align=PP_ALIGN.CENTER)
        for j, pt in enumerate(points):
            add_textbox(slide, f"  {pt}",
                        sx + Inches(0.2), Inches(2.1) + j * Inches(0.62),
                        sw - Inches(0.4), Inches(0.55),
                        font_size=Pt(15), color=C_GRAY)

    # 四大重点关注
    add_tag_rect(slide, "四大重点关注",
                 Inches(0.4), Inches(6.0), Inches(2.2), Inches(0.38), C_ACCENT)
    kws = ["量能变化", "政策落地", "行业景气", "公司业绩"]
    for i, kw in enumerate(kws):
        kx = Inches(0.4) + i * Inches(3.2)
        add_rect(slide, kx, Inches(6.5), Inches(3.0), Inches(0.7), C_CARD)
        add_textbox(slide, kw, kx, Inches(6.55), Inches(3.0), Inches(0.55),
                    font_size=Pt(16), bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# Slide 10 — 风险提示
# ═══════════════════════════════════════════════
def slide_risk():
    slide = blank_slide()
    fill_bg(slide, C_BG)
    slide_header(slide, "08  风险提示", "五大风险 · 应对措施")

    risks = [
        ("市场风险", "震荡调整可能持续，指数回落压力",    "控制仓位，分散投资"),
        ("政策风险", "政策落地进度不及预期，预期落空",    "密切关注政策动向"),
        ("流动性风险","税期临近（3月15日），资金面紧张",  "保持适当现金储备"),
        ("外围风险", "海外市场波动加剧，情绪传导影响",    "关注外围市场变化"),
        ("个股风险", "高位股分化明显，部分个股深度回调",  "避免追高，严格止损"),
    ]
    rw = Inches(12.5)
    rh = Inches(0.95)
    for i, (title, desc, action) in enumerate(risks):
        ry = Inches(1.1) + i * (rh + Inches(0.12))
        add_rect(slide, Inches(0.4), ry, rw, rh, C_CARD)
        # 左侧序号色块
        num_color = [C_RED, C_ORANGE, C_ORANGE, C_ORANGE, C_RED][i]
        add_rect(slide, Inches(0.4), ry, Inches(0.55), rh, num_color)
        add_textbox(slide, f"0{i+1}", Inches(0.4), ry, Inches(0.55), rh,
                    font_size=Pt(18), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, title,
                    Inches(1.1), ry + Inches(0.05), Inches(2.2), Inches(0.42),
                    font_size=Pt(17), bold=True, color=num_color)
        add_textbox(slide, desc,
                    Inches(3.5), ry + Inches(0.05), Inches(5.5), Inches(0.42),
                    font_size=Pt(14), color=C_GRAY)
        # 右侧应对
        add_rect(slide, Inches(9.2), ry + Inches(0.18), Inches(3.5), Inches(0.55), C_BG2)
        add_textbox(slide, f"应对：{action}",
                    Inches(9.3), ry + Inches(0.2), Inches(3.3), Inches(0.5),
                    font_size=Pt(13), color=C_ACCENT)

    # 底部提醒
    add_rect(slide, Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.7), C_CARD)
    add_rect(slide, Inches(0.4), Inches(6.5), Inches(0.08), Inches(0.7), C_RED)
    add_textbox(slide,
        "特别提醒：请严格执行止损纪律，单一持仓不超30%，保留不低于20%现金仓位！",
        Inches(0.6), Inches(6.55), Inches(12.2), Inches(0.55),
        font_size=Pt(15), bold=True, color=C_ORANGE)

# ═══════════════════════════════════════════════
# Slide 11 — 总结
# ═══════════════════════════════════════════════
def slide_summary():
    slide = blank_slide()
    fill_bg(slide, C_BG)
    slide_header(slide, "09  总结", "核心观点 · 数据回顾 · 投资建议")

    # 左侧：核心数据
    add_tag_rect(slide, "关键数据回顾", Inches(0.4), Inches(1.0),
                 Inches(2.4), Inches(0.38), C_ACCENT)
    data_lines = [
        ("上证指数",   "3366.16",  "-0.19%", C_RED),
        ("深证成指",   "10825.70", "-0.17%", C_RED),
        ("创业板指",   "2199.88",  "-0.25%", C_RED),
        ("成交额",     "15057亿",  "缩量3176亿", C_ORANGE),
        ("上涨家数",   "3279",     "+62.7%占比", C_GREEN),
        ("主力净流出", "248.13亿", "大幅流出", C_RED),
        ("北向净买入", "+50亿",    "逆势流入", C_GREEN),
    ]
    dw = Inches(6.0)
    dh = Inches(0.64)
    for i, (label, val, note, color) in enumerate(data_lines):
        dy = Inches(1.5) + i * dh
        add_rect(slide, Inches(0.4), dy, dw, dh - Inches(0.05), C_CARD)
        add_textbox(slide, label, Inches(0.5), dy, Inches(2.2), dh,
                    font_size=Pt(14), color=C_GRAY)
        add_textbox(slide, val, Inches(2.9), dy, Inches(1.8), dh,
                    font_size=Pt(15), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, note, Inches(4.8), dy, Inches(1.5), dh,
                    font_size=Pt(13), color=color, align=PP_ALIGN.RIGHT)

    # 右侧：核心结论
    add_tag_rect(slide, "核心结论", Inches(6.8), Inches(1.0),
                 Inches(2.0), Inches(0.38), C_ORANGE)
    conclusions = [
        ("震荡整理，缩量观望", C_ORANGE, "市场特征"),
        ("结构分化，北向逆势", C_ACCENT, "资金格局"),
        ("医药资源，把握轮动", C_GREEN,  "板块机会"),
        ("短期谨慎，中长期乐观", C_GREEN,"操作策略"),
    ]
    for i, (text, color, cat) in enumerate(conclusions):
        cy = Inches(1.5) + i * Inches(1.1)
        add_rect(slide, Inches(6.8), cy, Inches(6.2), Inches(0.95), C_CARD)
        add_rect(slide, Inches(6.8), cy, Inches(0.06), Inches(0.95), color)
        add_textbox(slide, cat, Inches(7.0), cy + Inches(0.06),
                    Inches(1.8), Inches(0.35), font_size=Pt(12), color=color)
        add_textbox(slide, text, Inches(7.0), cy + Inches(0.42),
                    Inches(5.8), Inches(0.42), font_size=Pt(16), bold=True, color=C_WHITE)

    # 底部免责
    add_rect(slide, Inches(0.4), Inches(6.8), Inches(12.5), Inches(0.5), C_CARD)
    add_textbox(slide,
        "免责声明：本报告仅供参考，不构成投资建议。股市有风险，投资需谨慎，请根据自身情况独立决策。",
        Inches(0.6), Inches(6.83), Inches(12.2), Inches(0.42),
        font_size=Pt(12), color=C_GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# Slide 12 — 封底
# ═══════════════════════════════════════════════
def slide_end():
    slide = blank_slide()
    fill_bg(slide, C_BG)

    add_rect(slide, 0, 0, Inches(0.2), H, C_ACCENT)
    add_rect(slide, 0, Inches(6.8), W, Inches(0.7), C_CARD)

    add_textbox(slide, "感谢观看",
                Inches(1), Inches(2.2), Inches(11), Inches(1.4),
                font_size=Pt(60), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "THANK YOU",
                Inches(1), Inches(3.5), Inches(11), Inches(0.8),
                font_size=Pt(28), color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, "2025年3月10日  |  WorkBuddy AI  |  数据来源：新浪财经 / 东方财富 / 雪球",
                Inches(1), Inches(6.88), Inches(11.333), Inches(0.45),
                font_size=Pt(14), color=C_GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# 构建所有幻灯片
# ═══════════════════════════════════════════════
print("Building slides...")
slide_cover()
print("  [1/12] Cover done")
slide_toc()
print("  [2/12] TOC done")
slide_overview()
print("  [3/12] Overview done")
slide_capital()
print("  [4/12] Capital done")
slide_sectors()
print("  [5/12] Sectors done")
slide_top10()
print("  [6/12] Top10 done")
slide_analysis()
print("  [7/12] Analysis done")
slide_outlook()
print("  [8/12] Outlook done")
slide_strategy()
print("  [9/12] Strategy done")
slide_risk()
print("  [10/12] Risk done")
slide_summary()
print("  [11/12] Summary done")
slide_end()
print("  [12/12] End done")

OUT = "D:/github/373Kice.github.io/2025年3月10日A股大盘汇报PPT.pptx"
prs.save(OUT)
print(f"\n[OK] PPT saved: {OUT}")
print(f"[OK] Total slides: {len(prs.slides)}")
