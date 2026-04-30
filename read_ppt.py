#!/usr/bin/env python
# -*- coding: utf-8 -*-
import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation

prs = Presentation('D:/github/373Kice.github.io/2025年3月10日A股大盘分析报告.pptx')
for i, slide in enumerate(prs.slides):
    print(f'=== Slide {i+1} ===')
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            print(shape.text[:500])
    print()
